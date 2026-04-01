"""
services/document_service.py
=============================
Document ingestion service — loads all supported file types from
``settings.docs_path`` into LangChain ``Document`` objects.

Supported formats
-----------------
- PDF  → PyPDFLoader
- CSV  → row-per-record with labeled columns
- PNG / JPG / JPEG → OCR via pytesseract

Usage::

    from services.document_service import document_service

    docs = document_service.load_all()          # fresh load
    docs = document_service.get_documents()     # cached load (loads once)
"""

from __future__ import annotations

import os
from pathlib import Path
from typing import Optional

from core.logging_config import get_logger
from configs.settings import settings

log = get_logger(__name__)


class DocumentService:
    """Loads and caches documents from the configured docs directory."""

    def __init__(self) -> None:
        self._documents: Optional[list] = None

    # ── public API ─────────────────────────────────────────────────────────

    def get_documents(self) -> list:
        """Return cached documents, loading them on first call."""
        if self._documents is None:
            self._documents = self.load_all()
        return self._documents

    def invalidate(self) -> None:
        """Clear cached documents so the next call reloads from disk."""
        self._documents = None

    def load_all(self) -> list:
        """Load all documents from ``settings.docs_path``. No caching."""
        docs_path: Path = settings.docs_path
        if not docs_path.exists():
            log.warning("Docs path does not exist: %s", docs_path)
            return []

        all_docs: list = []
        for fpath in sorted(docs_path.iterdir()):
            if not fpath.is_file():
                continue
            loaded = self._load_file(fpath)
            all_docs.extend(loaded)

        log.info("Loaded %d document chunk(s) from %s", len(all_docs), docs_path)
        return all_docs

    # ── per-file loaders ───────────────────────────────────────────────────

    def _load_file(self, fpath: Path) -> list:
        suffix = fpath.suffix.lower()
        try:
            if suffix == ".pdf":
                return self._load_pdf(fpath)
            if suffix == ".csv":
                return self._load_csv(fpath)
            if suffix in {".png", ".jpg", ".jpeg", ".webp"}:
                return self._load_image(fpath)
            if suffix == ".docx":
                return self._load_docx(fpath)
            if suffix == ".pptx":
                return self._load_pptx(fpath)
            if suffix == ".xlsx":
                return self._load_xlsx(fpath)
            if suffix == ".json":
                return self._load_json(fpath)
            if suffix in {".txt", ".md"}:
                return self._load_text(fpath)
        except Exception as exc:
            log.warning("Could not load %s: %s", fpath.name, exc)
        return []

    def _load_pdf(self, fpath: Path) -> list:
        from langchain_community.document_loaders import PyPDFLoader
        loader = PyPDFLoader(str(fpath))
        docs = []
        for doc in loader.load():
            doc.metadata["source"] = fpath.name
            docs.append(doc)
        log.debug("PDF loaded: %s (%d pages)", fpath.name, len(docs))
        return docs

    def _load_csv(self, fpath: Path) -> list:
        import pandas as pd
        from langchain_core.documents import Document

        df = pd.read_csv(fpath)
        docs = []
        for _, row in df.iterrows():
            row_text = ", ".join(
                f"{col}: {row[col]}" for col in df.columns
            )
            docs.append(
                Document(page_content=row_text, metadata={"source": fpath.name})
            )
        log.debug("CSV loaded: %s (%d rows)", fpath.name, len(docs))
        return docs

    def _load_image(self, fpath: Path) -> list:
        """Load an image via OCR with preprocessing for better text extraction.

        Preprocessing: grayscale → upscale (if small) → contrast boost.
        This mirrors the same pipeline in document_indexer_service to ensure
        that the project vector store and the windows-docs store both see the
        same OCR output for the same image file.
        """
        from langchain_core.documents import Document

        try:
            import pytesseract
            from PIL import Image, ImageEnhance
        except Exception as exc:
            log.warning("OCR dependency missing for %s: %s", fpath.name, exc)
            return []

        pytesseract.pytesseract.tesseract_cmd = settings.tesseract_cmd
        try:
            img = Image.open(fpath).convert("L")  # grayscale

            max_side = max(img.width, img.height)
            if max_side < 1600:
                scale = max(2, 1600 // max_side)
                img = img.resize(
                    (img.width * scale, img.height * scale),
                    Image.LANCZOS,
                )

            img = ImageEnhance.Contrast(img).enhance(2.0)
            extracted = pytesseract.image_to_string(img, config="--psm 6").strip()
        except Exception as exc:
            log.warning("OCR failed for %s: %s", fpath.name, exc)
            return []

        if extracted:
            log.debug("Image OCR loaded: %s", fpath.name)
            return [Document(page_content=extracted, metadata={"source": fpath.name})]
        else:
            log.info("No OCR text extracted from %s; skipped", fpath.name)
            return []

    def _load_docx(self, fpath: Path) -> list:
        import docx
        from langchain_core.documents import Document
        doc = docx.Document(str(fpath))
        text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        if text:
            log.debug("DOCX loaded: %s", fpath.name)
            return [Document(page_content=text, metadata={"source": fpath.name})]
        return []

    def _load_pptx(self, fpath: Path) -> list:
        from pptx import Presentation
        from langchain_core.documents import Document
        prs = Presentation(str(fpath))
        parts: list[str] = []
        for i, slide in enumerate(prs.slides, start=1):
            texts = [
                shape.text.strip()
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text.strip()
            ]
            if texts:
                parts.append(f"[Slide {i}]\n" + "\n".join(texts))
        text = "\n\n".join(parts)
        if text:
            log.debug("PPTX loaded: %s (%d slides)", fpath.name, len(prs.slides))
            return [Document(page_content=text, metadata={"source": fpath.name})]
        return []

    def _load_xlsx(self, fpath: Path) -> list:
        import pandas as pd
        from langchain_core.documents import Document
        xl = pd.ExcelFile(str(fpath))
        parts: list[str] = []
        for sheet in xl.sheet_names:
            df = xl.parse(sheet)
            if df.empty:
                continue
            lines = [
                ", ".join(f"{col}: {val}" for col, val in row.items())
                for _, row in df.iterrows()
            ]
            parts.append(f"[Sheet: {sheet}]\n" + "\n".join(lines))
        text = "\n\n".join(parts)
        if text:
            log.debug("XLSX loaded: %s", fpath.name)
            return [Document(page_content=text, metadata={"source": fpath.name})]
        return []

    def _load_json(self, fpath: Path) -> list:
        import json
        from langchain_core.documents import Document
        try:
            data = json.loads(fpath.read_text(encoding="utf-8", errors="ignore"))
            text = json.dumps(data, indent=2, ensure_ascii=False)
            if text:
                log.debug("JSON loaded: %s", fpath.name)
                return [Document(page_content=text, metadata={"source": fpath.name})]
        except Exception as exc:
            log.warning("JSON load error [%s]: %s", fpath.name, exc)
        return []

    def _load_text(self, fpath: Path) -> list:
        from langchain_core.documents import Document
        text = fpath.read_text(encoding="utf-8", errors="ignore").strip()
        if text:
            log.debug("Text loaded: %s", fpath.name)
            return [Document(page_content=text, metadata={"source": fpath.name})]
        return []


# ---------------------------------------------------------------------------
# Global singleton
# ---------------------------------------------------------------------------
document_service = DocumentService()
