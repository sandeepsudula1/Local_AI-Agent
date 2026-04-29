"""
services/document_indexer_service.py
=====================================
Scans configured document folders (recursively) for supported
file types, extracts text, splits into chunks, embeds, and persists in a
dedicated ChromaDB store separate from the project document store.

File-size tiers
---------------
* < 10 MB   → full chunking (all chunks stored)
* 10–100 MB → first ``MAX_CHUNKS_MEDIUM`` = 100 chunks only
* ≥ 100 MB  → metadata placeholder only (no text indexed)

Supported file types
--------------------
``.pdf``, ``.docx``, ``.txt``, ``.md``, ``.csv``, ``.png``, ``.jpg``, ``.jpeg``

Chunk parameters
----------------
chunk_size    ≈ 500 tokens  (2 000 characters at ~4 chars/token)
chunk_overlap ≈ 100 tokens  (  400 characters)

Change detection
----------------
``data/win_docs_index_state.json`` persists ``{absolute_path: mtime}`` for
every indexed file.  On each scan only new or modified files are re-indexed;
stale embeddings for modified files are deleted before new ones are added.

Integration
-----------
``services/vector_store_service.py`` calls ``document_indexer_service.start()``
after the project store is ready.  When this service finishes it re-registers
the ``documents.search`` tool so that queries search **both** stores.

Usage::

    from services.document_indexer_service import document_indexer_service

    document_indexer_service.start()      # non-blocking; runs in daemon thread
    document_indexer_service.wait()       # optional: block until finished
    db = document_indexer_service.get_vector_db()
"""

from __future__ import annotations

import json
import os
import threading
from pathlib import Path
from typing import Optional

from core.logging_config import get_logger
from configs.settings import settings, PROJECT_ROOT, DATA_DIR

log = get_logger(__name__)

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

SUPPORTED_EXTENSIONS: frozenset[str] = frozenset(
    {
        ".pdf", ".docx", ".txt", ".md", ".csv",
        ".png", ".jpg", ".jpeg", ".webp",
        ".pptx", ".xlsx", ".json",
        # Code / script files
        ".js", ".ts", ".py", ".java",
    }
)

_MB = 1024 * 1024
SIZE_FULL_THRESHOLD  = 10  * _MB   # below → full chunking
SIZE_LIMIT_THRESHOLD = 100 * _MB   # below → limited chunks; at/above → metadata only

MAX_CHUNKS_MEDIUM    = 100         # chunk cap for 10–100 MB files
CHROMA_BATCH_SIZE    = 500         # max chunks per Chroma add_documents() call

# ~500 tokens × 4 chars/token = 2 000 chars; ~100 tokens overlap = 400 chars
WIN_DOCS_CHUNK_SIZE    = 2_000
WIN_DOCS_CHUNK_OVERLAP = 400

_STATE_FILE: Path = Path(DATA_DIR) / "win_docs_index_state.json"


# ---------------------------------------------------------------------------
# Service class
# ---------------------------------------------------------------------------

class DocumentIndexerService:
    """
    Background service that indexes the Windows Documents folder into its own
    ChromaDB collection, then merges it with the project retrieval tool.
    """

    def __init__(self) -> None:
        self._db = None
        self._ready   = False
        self._loading = False
        self._lock    = threading.Lock()
        self._thread: Optional[threading.Thread] = None

    # ── public API ──────────────────────────────────────────────────────────

    @property
    def is_ready(self) -> bool:
        return self._ready

    @property
    def is_loading(self) -> bool:
        return self._loading

    def get_vector_db(self):
        """Return the Chroma instance, or ``None`` if not yet ready."""
        with self._lock:
            return self._db

    def start(self) -> None:
        """Start background indexing (returns immediately)."""
        if self._loading:
            log.debug("DocumentIndexerService.start() already running — skipped")
            return
        self._thread = threading.Thread(
            target=self._run,
            daemon=True,
            name="win-docs-indexer",
        )
        self._thread.start()

    def wait(self, timeout: float = 600.0) -> bool:
        """Block until ready or *timeout* seconds elapse.  Returns ``True`` when ready."""
        if self._thread:
            self._thread.join(timeout=timeout)
        return self._ready

    def sync_documents(self) -> dict:
        """Synchronise the vector store with the current state of the documents folder.

        Compares filenames present on disk against those indexed in ChromaDB,
        then removes embeddings for deleted files and indexes any new files.

        Operates on *filenames* (basename, not full path) so it works correctly
        whether the store was built from the root folder or from sub-folders.

        Returns
        -------
        dict
            {
                "folder_count":   int,   # files found on disk
                "indexed_count":  int,   # unique filenames in ChromaDB
                "removed":        int,   # embeddings deleted for removed files
                "added":          int,   # new files successfully indexed
            }
        """
        from langchain_text_splitters import RecursiveCharacterTextSplitter

        scan_root: Path = settings.windows_docs_path
        db = self.get_vector_db()

        # ── 1. Files currently on disk ────────────────────────────────────
        folder_files: set[str] = set()
        if scan_root.exists():
            for fpath in self._scan_files(scan_root):
                folder_files.add(fpath.name)

        print(f"[SYNC] Files in folder: {len(folder_files)}")

        # ── 2. Filenames already indexed in ChromaDB ──────────────────────
        indexed_files: set[str] = set()
        if db is not None:
            try:
                raw = db.get(include=["metadatas"])
                for meta in (raw.get("metadatas") or []):
                    fname = meta.get("source") or meta.get("file_name") or ""
                    if fname:
                        indexed_files.add(fname)
            except Exception as exc:
                log.warning("[SYNC] Could not read ChromaDB metadata: %s", exc)

        print(f"[SYNC] Indexed files: {len(indexed_files)}")

        # ── 3. Set operations ─────────────────────────────────────────────
        deleted_files: set[str] = indexed_files - folder_files
        new_files:     set[str] = folder_files  - indexed_files

        # ── 4. Remove embeddings for deleted files ────────────────────────
        removed = 0
        if deleted_files and db is not None:
            for fname in sorted(deleted_files):
                try:
                    db._collection.delete(where={"source": fname})
                    removed += 1
                    log.info("[SYNC] Removed embeddings for deleted file: %s", fname)
                except Exception as exc:
                    log.warning(
                        "[SYNC] Could not remove embeddings for %s: %s", fname, exc
                    )
            if removed:
                try:
                    db.persist()
                except Exception:
                    pass  # no-op in Chroma >= 0.4
                # Keep the state file in sync so the next full scan is accurate
                state = self._load_state()
                state = {k: v for k, v in state.items()
                         if Path(k).name not in deleted_files}
                self._save_state(state)

        # ── 5. Index new files ────────────────────────────────────────────
        added = 0
        if new_files and scan_root.exists():
            splitter = RecursiveCharacterTextSplitter(
                chunk_size=WIN_DOCS_CHUNK_SIZE,
                chunk_overlap=WIN_DOCS_CHUNK_OVERLAP,
            )
            new_chunks: list = []
            new_state_entries: dict = {}

            for fpath in self._scan_files(scan_root):
                if fpath.name not in new_files:
                    continue
                try:
                    chunks = self._index_file(fpath, splitter)
                    if not chunks:
                        log.debug("[SYNC] No content extracted from: %s", fpath.name)
                        continue
                    new_chunks.extend(chunks)
                    new_state_entries[str(fpath)] = fpath.stat().st_mtime
                    added += 1
                    log.info(
                        "[SYNC] Queued %s → %d chunk(s)",
                        fpath.name, len(chunks),
                    )
                except Exception as exc:
                    log.warning("[SYNC] Could not index %s: %s", fpath.name, exc)

            if new_chunks:
                store_path = str(settings.windows_docs_vector_store_path)
                if db is None:
                    # Store did not exist yet — create it from scratch
                    from langchain_community.embeddings import HuggingFaceEmbeddings
                    emb = HuggingFaceEmbeddings(
                        model_name=settings.embedding_model,
                        model_kwargs={"device": settings.embedding_device},
                    )
                    db = self._create_store_batched(new_chunks, emb, store_path)
                    with self._lock:
                        self._db = db
                    self._ready = True
                    self._post_ready()
                else:
                    self._add_to_store_batched(db, new_chunks)
                    try:
                        db.persist()
                    except Exception:
                        pass

                # Update state file with newly indexed entries
                state = self._load_state()
                state.update(new_state_entries)
                self._save_state(state)

        print(f"[SYNC] Removed: {removed}")
        print(f"[SYNC] Added:   {added}")

        log.info(
            "[SYNC] Complete — folder: %d | indexed: %d | removed: %d | added: %d",
            len(folder_files), len(indexed_files), removed, added,
        )
        return {
            "folder_count":  len(folder_files),
            "indexed_count": len(indexed_files),
            "removed":       removed,
            "added":         added,
        }

    # ── internals ───────────────────────────────────────────────────────────

    def _run(self) -> None:
        self._loading = True
        try:
            self._build_or_update()
        except Exception as exc:
            log.error("DocumentIndexerService fatal error: %s", exc, exc_info=True)
        finally:
            self._loading = False
            # Always re-register documents.search after the indexer finishes,
            # whether it succeeded, failed, or found nothing to update.
            self._post_ready()

    def _build_or_update(self) -> None:
        from langchain_community.embeddings import HuggingFaceEmbeddings
        from langchain_community.vectorstores import Chroma
        from langchain_text_splitters import RecursiveCharacterTextSplitter
        import shutil

        scan_root: Path = settings.windows_docs_path
        store_path: str = str(settings.windows_docs_vector_store_path)
        chroma_db_file  = settings.windows_docs_vector_store_path / "chroma.sqlite3"

        log.info("Windows docs indexer starting — scanning: %s", scan_root)

        if not scan_root.exists():
            log.warning("Windows docs path does not exist: %s", scan_root)
            return

        emb = HuggingFaceEmbeddings(
            model_name=settings.embedding_model,
            model_kwargs={"device": settings.embedding_device},
        )

        # ── Phase 1: load existing store IMMEDIATELY (before slow scan) ───────
        # We use chroma.sqlite3 as the authoritative existence check —
        # much more reliable than any(iterdir()) which can see temp files.
        settings.windows_docs_vector_store_path.mkdir(parents=True, exist_ok=True)
        db: object = None
        if chroma_db_file.exists():
            try:
                db = Chroma(persist_directory=store_path, embedding_function=emb)
                log.info("Windows docs store loaded from disk: %s", store_path)
                # Make the store available for queries RIGHT NOW — the scan
                # runs below and may take minutes; queries must not wait for it.
                with self._lock:
                    self._db = db
                self._ready = True
                log.info(
                    "Windows docs store ready (pre-scan). "
                    "post-scan will update if new files are found."
                )
            except Exception as exc:
                log.warning(
                    "Could not load Windows docs store (%s); will rebuild from scratch", exc
                )
                shutil.rmtree(store_path, ignore_errors=True)
                settings.windows_docs_vector_store_path.mkdir(parents=True, exist_ok=True)
                db = None

        # ── Phase 2: determine scan scope ─────────────────────────────────────
        subfolders = settings.windows_docs_subfolders  # tuple[str, ...]
        if subfolders:
            scan_roots = [
                scan_root / sub
                for sub in subfolders
                if (scan_root / sub).is_dir()
            ]
            if not scan_roots:
                log.warning(
                    "None of the configured WINDOWS_DOCS_SUBFOLDERS exist under %s; "
                    "falling back to full scan",
                    scan_root,
                )
                scan_roots = [scan_root]
            else:
                log.info(
                    "Restricting scan to %d configured subfolder(s): %s",
                    len(scan_roots),
                    [str(r) for r in scan_roots],
                )
        else:
            scan_roots = [scan_root]

        # ── Phase 3: incremental scan ─────────────────────────────────────────
        state    = self._load_state()
        splitter = RecursiveCharacterTextSplitter(
            chunk_size=WIN_DOCS_CHUNK_SIZE,
            chunk_overlap=WIN_DOCS_CHUNK_OVERLAP,
        )

        new_chunks:     list = []
        modified_pairs: list = []
        new_state:      dict = {}

        files_new = files_modified = files_skipped = 0

        for root in scan_roots:
            for fpath in self._scan_files(root):
                try:
                    mtime = fpath.stat().st_mtime
                    key   = str(fpath)
                    new_state[key] = mtime

                    if state.get(key) == mtime:
                        files_skipped += 1
                        continue

                    chunks = self._index_file(fpath, splitter)
                    if not chunks:
                        log.debug("No indexable content in: %s", fpath.name)
                        continue

                    if db is not None and key in state:
                        modified_pairs.append((fpath, chunks))
                        files_modified += 1
                    else:
                        new_chunks.extend(chunks)
                        files_new += 1

                    log.info(
                        "Indexed %s → %d chunk(s) (%.1f KB)",
                        fpath.name, len(chunks), fpath.stat().st_size / 1024,
                    )

                except Exception as exc:
                    log.warning("Skipping %s: %s", fpath, exc)

        total_scanned = files_new + files_modified + files_skipped
        log.info(
            "Windows docs scan complete: %d file(s) found — "
            "%d new, %d modified, %d unchanged",
            total_scanned, files_new, files_modified, files_skipped,
        )

        # ── remove embeddings for deleted files ───────────────────────────────
        files_deleted = 0
        if db is not None:
            deleted_paths = set(state.keys()) - set(new_state.keys())
            for deleted_path in deleted_paths:
                try:
                    db._collection.delete(where={"file_path": deleted_path})
                    files_deleted += 1
                    log.info(
                        "Removed embeddings for deleted file: %s",
                        Path(deleted_path).name,
                    )
                except Exception as exc:
                    log.debug(
                        "Could not remove embeddings for deleted file %s: %s",
                        deleted_path, exc,
                    )
            if files_deleted:
                log.info(
                    "Cleaned up embeddings for %d deleted file(s)", files_deleted
                )

        # ── remove stale embeddings for modified files ─────────────────────
        for fpath, chunks in modified_pairs:
            if db is not None:
                try:
                    db._collection.delete(where={"file_path": str(fpath)})
                    log.debug("Removed stale embeddings for: %s", fpath.name)
                except Exception as exc:
                    log.debug(
                        "Could not remove stale embeddings for %s: %s", fpath.name, exc
                    )
            new_chunks.extend(chunks)

        # ── persist new/modified chunks ────────────────────────────────────
        if not new_chunks:
            if db is None:
                log.warning("No indexable documents found under %s", scan_root)
            else:
                if files_deleted:
                    try:
                        db.persist()
                    except Exception:
                        pass
                self._save_state(new_state)
                if files_deleted:
                    log.info(
                        "[Ready] Windows docs store — removed %d deleted file(s)",
                        files_deleted,
                    )
                else:
                    log.info("[Ready] Windows docs store — no changes since last index")
            return

        if db is None:
            db = self._create_store_batched(new_chunks, emb, store_path)
        else:
            self._add_to_store_batched(db, new_chunks)

        try:
            db.persist()
        except Exception:
            pass  # no-op in Chroma >= 0.4

        self._save_state(new_state)

        with self._lock:
            self._db = db
        self._ready = True

        log.info(
            "[Ready] Windows docs updated: %d chunks added (%d new file(s), %d modified)",
            len(new_chunks), files_new, files_modified,
        )

    # ── batched Chroma helpers ─────────────────────────────────────────────────

    def _create_store_batched(self, chunks: list, emb, store_path: str):
        """Create a new Chroma store from *chunks* in safe-sized batches."""
        from langchain_community.vectorstores import Chroma
        db = None
        total = len(chunks)
        for i in range(0, total, CHROMA_BATCH_SIZE):
            batch = chunks[i : i + CHROMA_BATCH_SIZE]
            batch_num = i // CHROMA_BATCH_SIZE + 1
            total_batches = (total + CHROMA_BATCH_SIZE - 1) // CHROMA_BATCH_SIZE
            log.info(
                "Creating store — batch %d/%d (%d chunks)",
                batch_num, total_batches, len(batch),
            )
            if db is None:
                db = Chroma.from_documents(batch, emb, persist_directory=store_path)
            else:
                db.add_documents(batch)
        return db

    def _add_to_store_batched(self, db, chunks: list) -> None:
        """Add *chunks* to an existing Chroma store in safe-sized batches."""
        total = len(chunks)
        for i in range(0, total, CHROMA_BATCH_SIZE):
            batch = chunks[i : i + CHROMA_BATCH_SIZE]
            batch_num = i // CHROMA_BATCH_SIZE + 1
            total_batches = (total + CHROMA_BATCH_SIZE - 1) // CHROMA_BATCH_SIZE
            log.info(
                "Adding to store — batch %d/%d (%d chunks)",
                batch_num, total_batches, len(batch),
            )
            db.add_documents(batch)

    # ── file scanning ────────────────────────────────────────────────────────

    def _scan_files(self, root: Path):
        """Yield every supported file under *root*, skipping hidden / system dirs."""
        for dirpath, dirnames, filenames in os.walk(root):
            # Skip hidden directories (.git, .venv, …) and Windows system dirs ($RECYCLE.BIN, …)
            dirnames[:] = [
                d for d in dirnames
                if not d.startswith(".") and not d.startswith("$")
            ]
            dp = Path(dirpath)
            for fname in sorted(filenames):
                fpath = dp / fname
                if fpath.suffix.lower() in SUPPORTED_EXTENSIONS:
                    yield fpath

    # ── per-file indexing ────────────────────────────────────────────────────

    def _index_file(self, fpath: Path, splitter) -> list:
        """
        Return a list of LangChain ``Document`` chunks for *fpath*,
        applying file-size tier rules.
        """
        from langchain_core.documents import Document

        file_size = fpath.stat().st_size
        suffix    = fpath.suffix.lower()
        mtime     = fpath.stat().st_mtime

        base_metadata: dict = {
            "file_name":          fpath.name,
            "file_path":          str(fpath),
            "file_type":          suffix,
            "last_modified_time": mtime,
            "source":             fpath.name,
        }

        # ≥ 100 MB → store a metadata placeholder only; do not chunk content
        if file_size >= SIZE_LIMIT_THRESHOLD:
            log.info(
                "Large file (%.1f MB) — metadata placeholder only: %s",
                file_size / _MB, fpath.name,
            )
            return [Document(
                page_content=(
                    f"[File too large to index — content omitted. "
                    f"File: {fpath.name}, "
                    f"Size: {file_size / _MB:.1f} MB]"
                ),
                metadata={**base_metadata, "indexed": False, "reason": "file_too_large"},
            )]

        # Extract raw text
        text = self._extract_text(fpath, suffix)
        if not text or not text.strip():
            return []

        raw_doc = Document(page_content=text, metadata=base_metadata)
        chunks  = splitter.split_documents([raw_doc])

        # 10–100 MB → limit to first MAX_CHUNKS_MEDIUM chunks
        if file_size >= SIZE_FULL_THRESHOLD:
            log.info(
                "Medium file (%.1f MB) — capped at %d chunks: %s",
                file_size / _MB, MAX_CHUNKS_MEDIUM, fpath.name,
            )
            chunks = chunks[:MAX_CHUNKS_MEDIUM]

        return chunks

    # ── text extractors ──────────────────────────────────────────────────────

    def _extract_text(self, fpath: Path, suffix: str) -> str:
        if suffix == ".pdf":
            return self._read_pdf(fpath)
        if suffix == ".docx":
            return self._read_docx(fpath)
        if suffix == ".pptx":
            return self._read_pptx(fpath)
        if suffix in {".txt", ".md"}:
            return fpath.read_text(encoding="utf-8", errors="ignore")
        if suffix == ".csv":
            return self._read_csv(fpath)
        if suffix == ".xlsx":
            return self._read_xlsx(fpath)
        if suffix == ".json":
            return self._read_json(fpath)
        if suffix in {".png", ".jpg", ".jpeg", ".webp"}:
            return self._read_image_ocr(fpath)
        return ""

    def _read_pdf(self, fpath: Path) -> str:
        try:
            from langchain_community.document_loaders import PyPDFLoader
            pages = PyPDFLoader(str(fpath)).load()
            return "\n\n".join(p.page_content for p in pages)
        except Exception as exc:
            log.warning("PDF read error [%s]: %s", fpath.name, exc)
            return ""

    def _read_docx(self, fpath: Path) -> str:
        try:
            from docx import Document
            doc = Document(str(fpath))
            paragraphs = []
            for p in doc.paragraphs:
                text = p.text.strip()
                if text:
                    clean_text = " ".join(text.split())
                    paragraphs.append(clean_text)
            for table in doc.tables:
                for row in table.rows:
                    row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_text:
                        paragraphs.append(" | ".join([" ".join(c.split()) for c in row_text]))
            content = "\n".join(paragraphs)
            if content:
                print("[FILE] Clean text extracted from DOCX")
            return content
        except Exception as exc:
            log.warning("DOCX read error [%s]: %s", fpath.name, exc)
            return ""

    def _read_csv(self, fpath: Path) -> str:
        try:
            import pandas as pd
            df = pd.read_csv(fpath)
            lines = [
                ", ".join(f"{col}: {val}" for col, val in row.items())
                for _, row in df.iterrows()
            ]
            return "\n".join(lines)
        except Exception as exc:
            log.warning("CSV read error [%s]: %s", fpath.name, exc)
            return ""

    def _read_pptx(self, fpath: Path) -> str:
        try:
            from pptx import Presentation
            prs = Presentation(str(fpath))
            parts: list[str] = []
            for i, slide in enumerate(prs.slides, start=1):
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text.strip())
                if texts:
                    parts.append(f"[Slide {i}]\n" + "\n".join(texts))
            return "\n\n".join(parts)
        except Exception as exc:
            log.warning("PPTX read error [%s]: %s", fpath.name, exc)
            return ""

    def _read_xlsx(self, fpath: Path) -> str:
        try:
            import pandas as pd
            xl = pd.ExcelFile(str(fpath))
            parts: list[str] = []
            for sheet in xl.sheet_names:
                df = xl.parse(sheet)
                if df.empty:
                    continue
                lines = [
                    ", ".join(f"{col}: {val}" for col, val in row.items())
                    for _, row in df.iterrows()
                ]
                parts.append(f"[Sheet: {sheet}]\n" + "\n".join(lines))
            return "\n\n".join(parts)
        except Exception as exc:
            log.warning("XLSX read error [%s]: %s", fpath.name, exc)
            return ""

    def _read_json(self, fpath: Path) -> str:
        try:
            import json
            data = json.loads(fpath.read_text(encoding="utf-8", errors="ignore"))
            return json.dumps(data, indent=2, ensure_ascii=False)
        except Exception as exc:
            log.warning("JSON read error [%s]: %s", fpath.name, exc)
            return ""

    def _read_image_ocr(self, fpath: Path) -> str:
        """Extract text from an image using OCR (Tesseract).

        Preprocessing pipeline applied before OCR:
        1. Convert to grayscale (L mode) — reduces colour noise.
        2. Upscale images smaller than 1600 px on the longest side — Tesseract
           performs significantly better on high-resolution input.
        3. Enhance contrast — makes faint or low-contrast text readable.
        """
        try:
            import pytesseract
            from PIL import Image, ImageEnhance
        except Exception as exc:
            log.warning("Image OCR dependency error [%s]: %s", fpath.name, exc)
            return ""

        try:
            from core.runtime_paths import find_tesseract
            tesseract_cmd = find_tesseract()
            if tesseract_cmd:
                pytesseract.pytesseract.tesseract_cmd = tesseract_cmd

            img = Image.open(str(fpath)).convert("L")  # grayscale

            # Upscale so Tesseract has enough resolution to work with
            max_side = max(img.width, img.height)
            if max_side < 1600:
                scale = max(2, 1600 // max_side)
                img = img.resize(
                    (img.width * scale, img.height * scale),
                    Image.LANCZOS,
                )

            # Boost contrast (factor 2.0 works well for screenshots/documents)
            img = ImageEnhance.Contrast(img).enhance(2.0)

            # PSM 6 = assume a single uniform block of text (good for screenshots)
            text = pytesseract.image_to_string(img, config="--psm 6").strip()
            if not text:
                log.info("Image OCR found no readable text: %s", fpath.name)
            return text
        except Exception as exc:
            log.warning("Image OCR read error [%s]: %s", fpath.name, exc)
            return ""

    # ── state management ─────────────────────────────────────────────────────

    def _load_state(self) -> dict:
        """Return the persisted {path: mtime} index-state dict (or {} if absent)."""
        if _STATE_FILE.exists():
            try:
                return json.loads(_STATE_FILE.read_text(encoding="utf-8"))
            except Exception:
                pass
        return {}

    def _save_state(self, state: dict) -> None:
        """Persist the {path: mtime} index-state dict to disk."""
        _STATE_FILE.parent.mkdir(parents=True, exist_ok=True)
        _STATE_FILE.write_text(json.dumps(state, indent=2), encoding="utf-8")

    # ── on-demand folder indexing ─────────────────────────────────────────────

    def is_folder_indexed(self, folder_path: str) -> bool:
        """Return ``True`` if at least one file from *folder_path* is present in the index state.

        This is a lightweight disk-state check — it does NOT query the vector store.
        """
        state = self._load_state()
        norm = os.path.normcase(os.path.normpath(folder_path))
        norm_sep = norm + os.sep
        return any(
            os.path.normcase(os.path.normpath(k)).startswith(norm_sep)
            for k in state
        )

    def index_folder(self, folder_path: str, wait: bool = True, timeout: float = 120.0) -> bool:
        """Index all supported documents under *folder_path* into the existing store.

        Called when the user dynamically grants access to a new folder at runtime.
        If the store is not yet ready, the folder is queued and indexed as part of
        the normal startup scan instead.

        Parameters
        ----------
        folder_path:
            Absolute path to the folder to index.
        wait:
            When ``True`` (default) block until indexing completes or *timeout*
            seconds elapse.  When ``False``, run in a background daemon thread.
        timeout:
            Maximum seconds to wait when *wait* is ``True``.

        Returns
        -------
        bool
            ``True`` when at least one document was successfully indexed or was
            already indexed, ``False`` when zero documents could be found or
            indexed.
        """
        result_holder: list[bool] = []

        def _do_index() -> None:
            from langchain_community.embeddings import HuggingFaceEmbeddings
            from langchain_community.vectorstores import Chroma
            from langchain_text_splitters import RecursiveCharacterTextSplitter

            root = Path(folder_path)
            if not root.exists():
                log.warning("index_folder: path does not exist: %s", folder_path)
                result_holder.append(False)
                return

            log.info("index_folder: scanning %s", folder_path)

            splitter = RecursiveCharacterTextSplitter(
                chunk_size=WIN_DOCS_CHUNK_SIZE,
                chunk_overlap=WIN_DOCS_CHUNK_OVERLAP,
            )

            state = self._load_state()
            new_chunks: list = []
            new_state_entries: dict = {}

            for fpath in self._scan_files(root):
                try:
                    mtime = fpath.stat().st_mtime
                    key = str(fpath)
                    # Skip already-indexed, unchanged files
                    if state.get(key) == mtime:
                        log.debug("index_folder: already indexed, skipping %s", fpath.name)
                        continue
                    chunks = self._index_file(fpath, splitter)
                    if not chunks:
                        continue
                    new_chunks.extend(chunks)
                    new_state_entries[key] = mtime
                    log.info("index_folder: queued %s → %d chunk(s)", fpath.name, len(chunks))
                except Exception as exc:
                    log.warning("index_folder: skipping %s: %s", fpath, exc)

            if not new_chunks:
                # Nothing new to index — folder may already be indexed or be empty
                already_indexed = any(str(fpath) in state for fpath in self._scan_files(root))
                log.info(
                    "index_folder: no new documents in %s (already indexed: %s)",
                    folder_path, already_indexed,
                )
                # Ensure retrieval tool is registered even when there are no new chunks.
                # This handles the case where the startup Windows-docs indexer is still
                # running when the user grants access to a folder that was indexed in a
                # previous session (already_indexed=True but win_db not yet in coordinator).
                if already_indexed and self.get_vector_db() is not None:
                    log.info(
                        "[INDEX] Folder %s already indexed — refreshing retrieval registration",
                        folder_path,
                    )
                    self._post_ready()
                result_holder.append(already_indexed)
                return

            db = self.get_vector_db()
            store_path = str(settings.windows_docs_vector_store_path)

            try:
                if db is None:
                    emb = HuggingFaceEmbeddings(
                        model_name=settings.embedding_model,
                        model_kwargs={"device": settings.embedding_device},
                    )
                    db = self._create_store_batched(new_chunks, emb, store_path)
                    with self._lock:
                        self._db = db
                    self._ready = True
                else:
                    self._add_to_store_batched(db, new_chunks)
                    try:
                        db.persist()
                    except Exception:
                        pass

                # Persist updated state
                state.update(new_state_entries)
                self._save_state(state)

                log.info(
                    "index_folder: indexed %d chunk(s) from %d file(s) in %s",
                    len(new_chunks), len(new_state_entries), folder_path,
                )
                self._post_ready()

                # Post-index validation: confirm files are in index state and retrieval is ready
                _validated_count = sum(
                    1 for k in self._load_state()
                    if os.path.normcase(os.path.normpath(k)).startswith(
                        os.path.normcase(os.path.normpath(folder_path)) + os.sep
                    )
                )
                log.info(
                    "[INDEX] Folder ready for retrieval: %s  (%d file(s) in index state, %d chunk(s) added)",
                    folder_path, _validated_count, len(new_chunks),
                )
                result_holder.append(True)
            except Exception as exc:
                log.error("index_folder: failed to add documents from %s: %s", folder_path, exc)
                result_holder.append(False)

        t = threading.Thread(target=_do_index, daemon=True, name="folder-indexer")
        t.start()
        if wait:
            t.join(timeout=timeout)
        return bool(result_holder and result_holder[0])

    # ── post-ready hook ───────────────────────────────────────────────────────

    def _post_ready(self) -> None:
        """Notify the store coordinator of the Windows docs store.

        The coordinator re-registers ``documents.search`` with both the
        project store and this Windows store.  If the Windows store is
        not available (indexing failed), the project store is still used.
        """
        try:
            from core.tool_registry import update_retrieval_stores
            win_db = self.get_vector_db()
            if win_db is None:
                log.warning(
                    "Windows docs store not available after indexer finished; "
                    "retrieval will use project store only"
                )
            else:
                log.info(
                    "Windows docs store ready — notifying coordinator "
                    "(persist_directory=%s)",
                    getattr(win_db, "_persist_directory", "unknown"),
                )
            update_retrieval_stores(win_db=win_db)
        except Exception as exc:
            log.warning("Could not notify coordinator of Windows docs store: %s", exc)


# ---------------------------------------------------------------------------
# Module-level singleton
# ---------------------------------------------------------------------------

document_indexer_service = DocumentIndexerService()
