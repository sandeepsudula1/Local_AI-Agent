from __future__ import annotations

import os
import re
from typing import Optional, Tuple

import pandas as pd

from core.logging_config import get_logger

log = get_logger(__name__)

try:
    import ollama
    HAVE_OLLAMA = True
except Exception:
    ollama = None
    HAVE_OLLAMA = False

from configs.llm_config import MODEL
print(f"[LLM] Using model: {MODEL}")

ROOT = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
DOCS_PATH = os.path.join(ROOT, "data", "documents")


def _normalize_path(path: str) -> str:
    """Normalize a path for comparison: lowercase, forward slashes, no trailing slash."""
    return path.lower().replace("\\", "/").strip("/")


def _list_doc_files():
    if not os.path.exists(DOCS_PATH):
        return []
    return os.listdir(DOCS_PATH)


def _deduplicate_lines(text: str) -> str:
    """Remove duplicate lines (case-insensitive) from text while preserving order."""
    if not text:
        return text
    seen: set = set()
    out = []
    for line in text.splitlines():
        key = line.strip().lower()
        if key not in seen:
            seen.add(key)
            out.append(line)
    return "\n".join(out)


def _relevant_excerpt(text: str, query: str, max_chars: int = 4000) -> str:
    """Return up to max_chars of text starting from the most query-relevant position."""
    if len(text) <= max_chars:
        return text
    _STOP = {
        "what", "when", "where", "which", "that", "this", "have", "from", "with",
        "does", "about", "were", "will", "into", "been", "they", "them", "said",
        "tell", "show", "give", "know", "more", "some", "like", "also", "used",
        "then", "than", "there", "their", "these",
    }
    q_words = {w for w in re.findall(r"[a-z]+", query.lower()) if len(w) > 3} - _STOP
    if not q_words:
        return text[:max_chars]
    step = max(100, max_chars // 20)
    best_pos, best_score = 0, 0
    for pos in range(0, max(1, len(text) - max_chars + 1), step):
        chunk = text[pos:pos + max_chars].lower()
        score = sum(chunk.count(w) for w in q_words)
        if score > best_score:
            best_score = score
            best_pos = pos
    return text[best_pos:best_pos + max_chars]


# Stopwords excluded from the relevance keyword set
_KEYWORD_STOPWORDS: frozenset = frozenset({
    "what", "when", "where", "which", "that", "this", "have", "from", "with",
    "does", "about", "were", "will", "into", "been", "they", "them", "said",
    "tell", "show", "give", "know", "more", "some", "like", "also", "used",
    "then", "than", "there", "their", "these", "file", "docs", "document",
    "please", "could", "would", "should", "find", "look", "search", "help",
    "here", "those", "just", "only", "very", "much", "many",
    "to", "do", "the",
})

def answer_from_file(query: str, content: str, model_name: str = MODEL, file_path_used: str = None, is_summary: bool = False) -> str:
    """Answer a query using only the provided document content (Active File mode)."""
    try:
        if not content or not content.strip():
            return "The document is empty."
        
        source = os.path.basename(file_path_used) if file_path_used else "document"
        
        # Use standard relevance excerpt if too long
        if len(content) > 12000:
            context = _relevant_excerpt(content, query, max_chars=12000)
        else:
            context = content
            
        grounded_context = (
            "Answer ONLY using the content below.\n"
            "Do NOT say you cannot access files.\n\n"
            "DOCUMENT:\n"
            "----------------\n"
            f"{context}\n"
            "----------------"
        )
            
        return _ask_llm(model_name, grounded_context, query, source, is_summary=is_summary)
    except Exception as e:
        import traceback
        print(traceback.format_exc())
        print("[FILE ERROR]", str(e))
        return "Error processing file."


def _extract_query_keywords(query: str) -> set:
    """Return a set of meaningful lowercase tokens from *query* (min 3 chars, non-stopword)."""
    return {
        w for w in re.findall(r"[a-z0-9]+", query.lower())
        if len(w) >= 3 and w not in _KEYWORD_STOPWORDS
    }


def _keyword_filter(docs: list, query: str) -> list:
    """Return the subset of *docs* whose page_content contains at least one query keyword.

    When no keyword can be extracted from *query* (e.g. the query is too short
    or consists entirely of stopwords) the original list is returned unchanged so
    that normal retrieval still proceeds.

    An empty return list means "no chunks are relevant to this query" — the
    caller should respond with a 'not found' message rather than sending
    unrelated context to the LLM.
    """
    keywords = _extract_query_keywords(query)
    if not keywords:
        return docs  # cannot judge relevance — pass through
    return [
        doc for doc in docs
        if any(kw in doc.page_content.lower() for kw in keywords)
    ]


def _detect_target_file(query: str):
    """Return the specific filename the user is asking about, or None."""
    q = query.strip()
    available = _list_doc_files()

    # 1. Direct substring match against known filenames (case-insensitive)
    q_lower = q.lower()
    for fname in available:
        if fname.lower() in q_lower:
            return fname

    # 2. Regex for any explicit filename pattern in the query
    #    If the user typed an explicit "name.ext" the match fires even for files NOT in
    #    data/documents/ (they live in the win_docs store).  In that case we return None
    #    so that Rule-1b (vector metadata search) handles the lookup.  We deliberately
    #    do NOT fall through to the keyword match (step 3) because that would incorrectly
    #    return a completely different local file (e.g. sandeep_internship_work.pdf when
    #    the user asked about Sandeep_S_Java\ Fullstack.pdf).
    m = re.search(r"[\w\s\-\.,()]+\.(?:pdf|pptx|csv|txt|md|png|jpg|jpeg|webp|docx|xlsx|json)", q, flags=re.IGNORECASE)
    if m:
        candidate = m.group(0).strip()
        # Exact match
        for fname in available:
            if fname.lower() == candidate.lower():
                return fname
        # Stem exact match (not substring — 'sandeep' ≠ 'sandeep_internship_work')
        cand_stem = os.path.splitext(candidate)[0].lower().strip()
        for fname in available:
            if cand_stem == os.path.splitext(fname)[0].lower():
                return fname
        # An explicit filename was referenced but it is NOT in data/documents/.
        # Stop here — Rule-1b will search the win_docs vector store instead.
        return None

    # 3. Keyword match — only reaches here when the query contains NO explicit
    #    filename pattern, so we can safely keyword-match against local docs.
    #    e.g. "summarize the internship pdf" → matches sandeep_internship_work.pdf
    for fname in available:
        ext = os.path.splitext(fname)[1].lower()
        if ext not in {".pdf", ".csv", ".txt", ".md", ".png", ".jpg", ".jpeg", ".webp",
                       ".docx", ".pptx", ".xlsx", ".json"}:
            continue
        stem_words = re.findall(r"[a-z]+", os.path.splitext(fname)[0].lower())
        # Match if any meaningful stem word (len>3) appears in the query
        if any(w in q_lower for w in stem_words if len(w) > 3):
            return fname

    return None


def _load_file_content(fname: str):
    """Load text content from a document file. Returns (text, fname)."""
    fpath = os.path.join(DOCS_PATH, fname)
    if not os.path.exists(fpath):
        return None, None
    ext = os.path.splitext(fname)[1].lower()

    if ext == ".pdf":
        try:
            from langchain_community.document_loaders import PyPDFLoader
            docs = list(PyPDFLoader(fpath).load())
            return "\n\n".join(d.page_content for d in docs), fname
        except Exception as e:
            return f"(Error reading PDF: {e})", fname

    elif ext == ".csv":
        try:
            df = pd.read_csv(fpath)
            rows = []
            for i, row in df.iterrows():
                row_parts = [f"{col}={row[col]}" for col in df.columns]
                rows.append(f"Record {i + 1}: {', '.join(row_parts)}")
            table_desc = "\n".join(rows)
            return (
                f"CSV file: {fname}\n"
                f"Columns: {', '.join(df.columns.tolist())}\n"
                f"Number of records: {len(df)}\n"
                f"(Each record is data for the SAME entity, e.g. the same company across different years)\n\n"
                + table_desc
            ), fname
        except Exception as e:
            return f"(Error reading CSV: {e})", fname

    elif ext == ".txt":
        try:
            return open(fpath, "r", encoding="utf-8", errors="ignore").read(), fname
        except Exception as e:
            return f"(Error reading TXT: {e})", fname

    elif ext in {".png", ".jpg", ".jpeg", ".webp"}:
        try:
            import pytesseract
            from PIL import Image, ImageEnhance
            pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
            img = Image.open(fpath).convert("L")
            max_side = max(img.width, img.height)
            if max_side < 1600:
                scale = max(2, 1600 // max_side)
                img = img.resize((img.width * scale, img.height * scale), Image.LANCZOS)
            img = ImageEnhance.Contrast(img).enhance(2.0)
            text = pytesseract.image_to_string(img, config="--psm 6").strip()
            if text:
                return text, fname
            return None, fname   # no readable text — skip LLM, return helpful message below
        except Exception:
            return None, fname   # Tesseract not installed or failed

    elif ext == ".pptx":
        try:
            from pptx import Presentation
            prs = Presentation(fpath)
            parts = []
            for i, slide in enumerate(prs.slides, start=1):
                texts = [
                    shape.text.strip()
                    for shape in slide.shapes
                    if hasattr(shape, "text") and shape.text.strip()
                ]
                if texts:
                    parts.append(f"[Slide {i}]\n" + "\n".join(texts))
            return "\n\n".join(parts) or None, fname
        except Exception as e:
            return f"(Error reading PPTX: {e})", fname

    elif ext == ".xlsx":
        try:
            df_sheets = pd.read_excel(fpath, sheet_name=None)
            parts = []
            for sheet, df in df_sheets.items():
                if df.empty:
                    continue
                rows = []
                for i, row in df.iterrows():
                    row_parts = [f"{col}={row[col]}" for col in df.columns]
                    rows.append(f"Record {i + 1}: {', '.join(row_parts)}")
                parts.append(f"[Sheet: {sheet}]\n" + "\n".join(rows))
            return "\n\n".join(parts) or None, fname
        except Exception as e:
            return f"(Error reading XLSX: {e})", fname

    elif ext == ".json":
        try:
            import json
            data = json.loads(open(fpath, encoding="utf-8", errors="ignore").read())
            return json.dumps(data, indent=2, ensure_ascii=False), fname
        except Exception as e:
            return f"(Error reading JSON: {e})", fname

    elif ext == ".docx":
        try:
            import docx
            doc = docx.Document(fpath)
            return "\n".join(p.text for p in doc.paragraphs), fname
        except Exception as e:
            return f"(Error reading DOCX: {e})", fname

    return "(Unsupported file type)", fname


def _extract_filename_from_query(query: str) -> Optional[str]:
    """Extract a bare filename (with extension) from a query, if present.

    Unlike ``_detect_target_file`` this function does NOT require the file to
    exist on disk — it just applies a regex so it works for files indexed in
    the Windows docs vector store.
    """
    _EXTS = r'(?:pdf|pptx|docx|txt|md|csv|xlsx|xls|png|jpg|jpeg|py|js|ts|json|html|xml|ini|cfg|yaml|yml)'
    _WORD = r'[\w][\w\-\.]*'
    m = re.search(rf'\b({_WORD}\.{_EXTS})\b', query, re.IGNORECASE)
    return m.group(1).strip() if m else None


def _find_file_under_root(root_path: str, filename: str) -> Optional[str]:
    """Find a filename under a root directory (case-insensitive) and return full path.

    Matching strategy (in order):
    1. Exact case-insensitive match  (e.g. ``filename`` = "report.pdf")
    2. Suffix match  — the stored filename ends with " " + target, which handles
       partial names extracted from queries with spaces (e.g. ``filename`` =
       "140804.png" correctly finds "Screenshot 2025-10-28 140804.png").
    """
    if not root_path or not os.path.exists(root_path):
        return None
    target = filename.lower()
    for dirpath, _, filenames in os.walk(root_path):
        for fname in filenames:
            fname_lower = fname.lower()
            if fname_lower == target or fname_lower.endswith(" " + target):
                return os.path.join(dirpath, fname)
    return None


def _ocr_image_file(file_path: str) -> Optional[str]:
    """Extract OCR text from an image file with preprocessing.

    Preprocessing: grayscale → upscale (if < 1600 px on longest side) → contrast
    boost.  This produces more consistent, higher-quality OCR output compared to
    passing the raw image directly to Tesseract.
    Returns None if dependencies are unavailable or no text is extracted.
    """
    try:
        import pytesseract
        from PIL import Image, ImageEnhance
    except Exception:
        return None

    try:
        tesseract_path = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
        if os.path.exists(tesseract_path):
            pytesseract.pytesseract.tesseract_cmd = tesseract_path

        img = Image.open(file_path).convert("L")  # grayscale

        max_side = max(img.width, img.height)
        if max_side < 1600:
            scale = max(2, 1600 // max_side)
            img = img.resize(
                (img.width * scale, img.height * scale),
                Image.LANCZOS,
            )

        img = ImageEnhance.Contrast(img).enhance(2.0)
        text = pytesseract.image_to_string(img, config="--psm 6").strip()
        return text or None
    except Exception:
        return None


def _load_document_from_path(file_path: str) -> Optional[str]:
    """Load text content from a document at an arbitrary absolute file path.

    Supports the same formats as ``_load_file_content`` but works with any
    path rather than being restricted to ``DOCS_PATH``.  Used as a fallback
    when a file exists on disk but has not yet been indexed in the vector store.
    """
    if not os.path.exists(file_path):
        return None
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf":
        try:
            from langchain_community.document_loaders import PyPDFLoader
            docs = list(PyPDFLoader(file_path).load())
            text = "\n\n".join(d.page_content for d in docs).strip()
            if not text:
                return "Error: PDF is empty or contains only images."
            return text
        except Exception as e:
            return f"Error reading PDF: {e}"
    elif ext == ".txt":
        try:
            content = open(file_path, "r", encoding="utf-8", errors="ignore").read().strip()
            if content:
                log.debug("[LOAD] %r: %d chars (utf-8)", file_path, len(content))
                return content
            # Try latin-1 before giving up (handles some Windows-encoded files)
            content = open(file_path, "r", encoding="latin-1", errors="ignore").read().strip()
            if content:
                log.debug("[LOAD] %r: %d chars (latin-1)", file_path, len(content))
                return content
            log.warning("[LOAD] %r: file is empty", file_path)
            return None
        except Exception as exc:
            log.warning("[LOAD] Failed to read %r: %s", file_path, exc)
            return None
    elif ext in (".png", ".jpg", ".jpeg", ".webp"):
        return _ocr_image_file(file_path)
    elif ext == ".pptx":
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            parts = []
            for i, slide in enumerate(prs.slides, start=1):
                texts = [
                    s.text.strip() for s in slide.shapes
                    if hasattr(s, "text") and s.text.strip()
                ]
                if texts:
                    parts.append(f"[Slide {i}]\n" + "\n".join(texts))
            return "\n\n".join(parts) or None
        except Exception:
            return None
    elif ext in (".xlsx", ".xls"):
        try:
            df_sheets = pd.read_excel(file_path, sheet_name=None)
            parts = []
            for sheet, df in df_sheets.items():
                if df.empty:
                    continue
                rows = [
                    f"Record {i + 1}: {', '.join(f'{c}={r[c]}' for c in df.columns)}"
                    for i, r in df.iterrows()
                ]
                parts.append(f"[Sheet: {sheet}]\n" + "\n".join(rows))
            return "\n\n".join(parts) or None
        except Exception:
            return None
    elif ext == ".json":
        try:
            import json
            return json.dumps(
                json.loads(open(file_path, encoding="utf-8", errors="ignore").read()),
                indent=2, ensure_ascii=False,
            ) or None
        except Exception:
            return None
    elif ext == ".docx":
        try:
            import docx
            doc = docx.Document(file_path)
            text = "\n".join(p.text for p in doc.paragraphs).strip()
            if not text:
                return "Error: DOCX is empty."
            return text
        except Exception as e:
            return f"Error reading DOCX: {e}"
    elif ext == ".csv":
        try:
            df = pd.read_csv(file_path)
            rows = [
                f"Record {i + 1}: {', '.join(f'{c}={r[c]}' for c in df.columns)}"
                for i, r in df.iterrows()
            ]
            return (
                f"CSV: {os.path.basename(file_path)}\n"
                f"Columns: {', '.join(df.columns.tolist())}\n"
                f"Records: {len(df)}\n\n" + "\n".join(rows)
            ) or None
        except Exception:
            return None
    # General plain-text / code files — read as UTF-8 with error tolerance
    _TEXT_CODE_EXTS = {
        ".js", ".ts", ".py", ".java", ".html", ".xml",
        ".yaml", ".yml", ".ini", ".cfg", ".md", ".log",
        ".sh", ".bat", ".css", ".scss", ".sql",
    }
    if ext in _TEXT_CODE_EXTS:
        try:
            return open(file_path, "r", encoding="utf-8", errors="ignore").read() or None
        except Exception:
            return None
    return None


def _get_supported_files_in_folder(folder_path: str) -> list:
    """Return all supported document files in a folder.
    
    Returns list of (filepath, filename) tuples for supported formats.
    """
    if not os.path.isdir(folder_path):
        return []
    
    supported_exts = {
        ".pdf", ".csv", ".txt", ".md", ".png", ".jpg", ".jpeg", ".webp",
        ".docx", ".pptx", ".xlsx", ".json"
    }
    
    files = []
    try:
        for entry in os.scandir(folder_path):
            if entry.is_file():
                _, ext = os.path.splitext(entry.name)
                if ext.lower() in supported_exts:
                    files.append((entry.path, entry.name))
    except Exception:
        pass
    
    return sorted(files, key=lambda x: x[1])


def _load_all_files_from_folder(folder_path: str, query: str, model_name: str) -> Tuple[Optional[str], Optional[str]]:
    """Load and process all supported files from a folder.
    
    Returns combined content from all files in the folder.
    """
    files = _get_supported_files_in_folder(folder_path)
    
    if not files:
        return (
            f"No supported document files found in:\n📁 {folder_path}\n\n"
            f"Supported formats: PDF, TXT, CSV, MD, PNG, JPG, DOCX, PPTX, XLSX, JSON",
            None,
        )
    
    all_contents = []
    summary = _is_summary_intent(query)
    
    for file_path, file_name in files:
        log.info("Loading file from folder: %s", file_path)
        content = _load_document_from_path(file_path)
        if content:
            all_contents.append(f"=== {file_name} ===\n{content}")
        else:
            all_contents.append(f"=== {file_name} ===\n(Could not read file)")
    
    combined_content = "\n\n".join(all_contents)
    
    if not combined_content:
        file_list = ", ".join([f[1] for f in files])
        return (
            f"Found {len(files)} file(s) in {folder_path}:\n{file_list}\n\n"
            f"But could not read their contents.",
            folder_path,
        )
    
    # Ask LLM to process the content if needed
    if _is_summary_intent(query) or re.search(r'\b(read|show|list|get|summarize|explain)\b', query, re.IGNORECASE):
        answer = _ask_llm(
            model_name,
            combined_content if len(combined_content) <= 16000 else combined_content[:16000],
            query,
            os.path.basename(folder_path),
            is_summary=summary
        )
        if answer:
            return answer, folder_path
    
    # If no LLM answer or not summary intent, return formatted content
    preview = combined_content[:6000] if len(combined_content) > 6000 else combined_content
    truncation_msg = "\n\n...(truncated, showing first 6000 chars)" if len(combined_content) > 6000 else ""
    return (
        f"Files in {os.path.basename(folder_path)}:\n\n{preview}{truncation_msg}",
        folder_path,
    )
def _get_authorized_docs_root() -> str:
    """Return the authorized document root path (lowercase, normalised)."""
    try:
        from configs.settings import settings
        return str(settings.windows_docs_path).lower().rstrip("\\").rstrip("/")
    except Exception:
        return r"c:\ai_test_documents"


def _query_references_unauthorized_path(query: str) -> bool:
    """Return True when the query explicitly references a Windows path that is
    NOT under any authorized folder root (static ALLOWED_FOLDERS or dynamically
    granted folders).

    Uses path-boundary matching to avoid false negatives like treating
    'C:\\AI_Test_Documents2' as authorized just because it starts with
    'C:\\AI_Test_Documents'.
    """
    # Build the full list of authorized roots (static + dynamically granted)
    authorized_roots: list[str] = []
    try:
        from core.access_control import ALLOWED_FOLDERS as _AF
        authorized_roots = [p.lower().replace("/", "\\").rstrip("\\") for p in _AF]
    except Exception:
        authorized_roots = [_get_authorized_docs_root()]
    try:
        from core.permission_store import permission_store as _ps
        authorized_roots += [
            p.lower().replace("/", "\\").rstrip("\\")
            for p in _ps.get_granted_folders()
        ]
    except Exception:
        pass

    # Match any Windows-style path token: letter colon backslash (or forward slash)
    for m in re.finditer(r'[A-Za-z]:[/\\][^\s"\',;]*', query):
        found = m.group(0).lower().replace("/", "\\").rstrip("\\")
        # Path-boundary check: exact match OR properly nested under an authorized root
        if not any(
            found == auth or found.startswith(auth + "\\")
            for auth in authorized_roots
        ):
            return True
    return False


def _is_summary_intent(query: str) -> bool:
    """Return True when the user wants a summary, explanation, or fuller detail."""
    return bool(re.search(
        r'\b('
        r'summarize|summarise|summary|summarization'
        r'|explain|explanation'
        r'|describe|description'
        r'|overview|outline'
        r'|elaborate|expand'
        r'|full details?|all details?|more details?|give details'
        r'|tell me (more|about|everything|all)'
        r'|more (about|information|info|on)'
        r'|what (does|is|are|\w+ say|\w+ contain|\w+ cover|\w+ discuss)'
        r'|what is in|what.?s in|what.?s inside'
        r'|give me|show me everything|show me all'
        r'|key points?|main points?'
        r')\b',
        query.lower(),
    ))


def _search_by_filename_in_stores(filename: str, dbs: list) -> list:
    """Return ALL Document chunks whose metadata matches *filename*.

    Search order (stops at first hit per store):
    1. ``file_name`` exact match (ChromaDB ``where`` filter)
    2. ``file_name`` case-insensitive match (try lowercase variant)
    3. ``source`` exact / lowercase match
    4. Similarity search on the filename string, post-filtered by ``file_name``
       or ``source`` (handles stores with no operator support).
    """
    from langchain_core.documents import Document
    fname_lower = filename.lower()
    all_docs: list = []

    def _get_docs(db, where_filter: dict) -> list:
        """Query ChromaDB using a metadata WHERE filter.

        Uses the raw ``_collection.get()`` ChromaDB API which reliably supports
        WHERE filters on metadata.  LangChain's ``db.get()`` wrapper does NOT
        propagate the ``where`` kwarg to the underlying collection, so it always
        returned an empty result — that was the root cause of Rule-1b failures.
        """
        try:
            raw = db._collection.get(
                where=where_filter,
                include=["documents", "metadatas"],
            )
            docs_raw  = raw.get("documents") or []
            metas_raw = raw.get("metadatas") or [{}] * len(docs_raw)
            return [
                Document(page_content=c, metadata=m or {})
                for c, m in zip(docs_raw, metas_raw)
            ]
        except Exception as exc:
            log.debug("_get_docs WHERE filter failed %s: %s", where_filter, exc)
            return []

    for db in dbs:
        if db is None:
            continue

        # ─ attempt 1: file_name exact ────────────────────────────────────────
        docs = _get_docs(db, {"file_name": {"$eq": filename}})
        if not docs:
            # ─ attempt 1b: file_name lowercase ────────────────────────────
            docs = _get_docs(db, {"file_name": {"$eq": fname_lower}})
        if not docs:
            # ─ attempt 2: source field exact ──────────────────────────────
            docs = _get_docs(db, {"source": {"$eq": filename}})
        if not docs:
            docs = _get_docs(db, {"source": {"$eq": fname_lower}})

        if docs:
            all_docs.extend(docs)
            log.info(
                "_search_by_filename: metadata match — %d chunk(s) for %r in %s",
                len(docs), filename, getattr(db, "_persist_directory", "store"),
            )
            continue

        # ─ attempt 3: full metadata suffix scan ──────────────────────────────
        # Handles multi-word filenames such as 'Sandeep_S_Java Fullstack.pdf'
        # when the regex only extracted the last segment 'Fullstack.pdf'.
        # We fetch all chunks and filter by suffix instead of exact match.
        try:
            raw = db._collection.get(include=["documents", "metadatas"])
            docs_raw  = raw.get("documents") or []
            metas_raw = raw.get("metadatas") or [{}] * len(docs_raw)
            suffix_docs = [
                Document(page_content=c, metadata=m or {})
                for c, m in zip(docs_raw, metas_raw)
                if (
                    (m or {}).get("file_name", "").lower().endswith(" " + fname_lower)
                    or (m or {}).get("source", "").lower().endswith(" " + fname_lower)
                )
            ]
            if suffix_docs:
                all_docs.extend(suffix_docs)
                log.info(
                    "_search_by_filename: suffix scan — %d chunk(s) for %r",
                    len(suffix_docs), filename,
                )
                continue
        except Exception as exc:
            log.debug("Suffix metadata scan failed: %s", exc)

        # ─ attempt 4: similarity search + post-filter ─────────────────────
        # Post-filter uses exact match AND suffix match so that a partially
        # extracted name like 'Fullstack.pdf' still finds 'Sandeep_S_Java Fullstack.pdf'.
        try:
            hits = db.similarity_search_with_score(filename, k=200)
            matched = [
                doc for doc, _ in hits
                if (
                    doc.metadata.get("file_name", "").lower() == fname_lower
                    or doc.metadata.get("source",    "").lower() == fname_lower
                    # suffix match: 'fullstack.pdf' matches 'sandeep_s_java fullstack.pdf'
                    or doc.metadata.get("file_name", "").lower().endswith(" " + fname_lower)
                    or doc.metadata.get("source",    "").lower().endswith(" " + fname_lower)
                )
            ]
            if matched:
                all_docs.extend(matched)
                log.info(
                    "_search_by_filename: similarity fallback — %d chunk(s) for %r",
                    len(matched), filename,
                )
        except Exception as exc:
            log.debug("Filename search fallback failed: %s", exc)

    return all_docs


def _best_source_from_docs(docs: list, fallback: str) -> str:
    """Pick the most frequent metadata source/file_name from matched chunks."""
    counts: dict[str, int] = {}
    display: dict[str, str] = {}
    for d in docs:
        meta = d.metadata or {}
        src = (meta.get("file_name") or meta.get("source") or "").strip()
        if not src:
            continue
        key = src.lower()
        display.setdefault(key, src)
        counts[key] = counts.get(key, 0) + 1
    if not counts:
        return fallback
    best_key = max(counts, key=lambda k: counts[k])
    return display.get(best_key, fallback)


def _detect_content_type(context: str) -> str:
    """Infer the nature of the OCR/document content to choose the right prompt.

    Returns one of: ``"resume"``, ``"code_tree"``, ``"screenshot_ui"``, ``"generic"``.
    """
    if not context:
        return "generic"
    lower = context.lower()

    # File-tree / VS Code explorer screenshot
    file_tree_signals = sum([
        bool(re.search(r"\.(js|ts|py|css|html|json|jsx|tsx|vue|php|java|cs)\b", lower)),
        bool(re.search(r"\b(node_modules|src|components|services|public|dist|build|package\.json|\.gitignore)\b", lower)),
        bool(re.search(r"\bindex\.(js|ts|html|css)\b", lower)),
        context.count(".js") + context.count(".ts") + context.count(".py") > 3,
    ])
    if file_tree_signals >= 2:
        return "code_tree"

    # Resume / CV / profile
    resume_signals = sum([
        bool(re.search(r"\b(education|university|college|degree|bachelor|master|phd)\b", lower)),
        bool(re.search(r"\b(experience|internship|work experience|employment)\b", lower)),
        bool(re.search(r"\b(skills|certifications?|achievements?|awards?|publications?)\b", lower)),
        bool(re.search(r"\b(resume|curriculum vitae|cv\b|objective|summary)\b", lower)),
        bool(re.search(r"\b(gpa|cgpa|percentage|grade)\b", lower)),
    ])
    if resume_signals >= 2:
        return "resume"

    return "generic"


def _ask_llm(model_name, context, query, source, *, is_summary: bool = False):
    print(f"[PROMPT_SIZE] {len(context)} chars")
    print(f"[GENERATION] Starting grounded LLM generation")
    """Ask Ollama to answer using only the provided context. Returns answer string or None."""
    if not HAVE_OLLAMA:
        return None
    try:
        if is_summary:
            content_type = _detect_content_type(context)

            if content_type == "resume":
                system_prompt = (
                    "You are a helpful document summarization assistant. "
                    "Write a clear, faithful summary using ONLY the provided content. "
                    "Do not add facts that are not present in the content. "
                    "Structure the summary with short headings ONLY for sections that contain "
                    "actual information in the content: "
                    "Profile, Education, Internship/Experience, Technical Skills, Certifications, Projects. "
                    "IMPORTANT: If a section has no real content in the source, skip it entirely. "
                    "NEVER write 'None mentioned', 'Not provided', 'No X information', or any similar "
                    "placeholder — simply leave that heading out. "
                    "Do NOT say 'Based on the context' — write the summary directly."
                )
            elif content_type == "code_tree":
                system_prompt = (
                    "You are a helpful technical assistant. "
                    "The content below is OCR text extracted from a screenshot of a code editor or "
                    "file explorer panel. "
                    "Describe what the screenshot shows: the project name, folder structure, "
                    "source files, and what type of application it appears to be. "
                    "Use ONLY information visible in the content — do NOT invent sections like "
                    "'Education', 'Internship', 'Certifications', or any resume-like headings. "
                    "Do NOT say 'Based on the context' — describe it directly."
                )
            else:
                system_prompt = (
                    "You are a helpful document summarization assistant. "
                    "Write a clear, faithful summary or explanation using ONLY the provided content. "
                    "Do not add facts, sections, or headings that are not supported by the content. "
                    "NEVER write placeholder lines like 'None mentioned', 'Not provided', or "
                    "'No X information is given' — if information is absent, simply omit that topic. "
                    "Do NOT say 'Based on the context' — write directly."
                )

            summary_context = context if len(context) <= 16000 else context[:16000]
            user_prompt = f"""Summarize the following document:

{summary_context}
"""
            max_tokens = 1100
        else:
            system_prompt = (
                "You are a Local Multi-Agent AI Assistant running on the user's computer. "
                "You have access ONLY to documents indexed from the configured folder. "
                "Answer using ONLY facts explicitly stated in the CONTEXT below. "
                "If the user mentions a specific filename, answer ONLY from that file's content. "
                "Never use content from a different document than the one retrieved. "
                "State numbers, names, and dates exactly as they appear in the context. "
                "For CSV/tabular data: each ROW is a separate record for the SAME entity — "
                "do NOT treat different rows as different companies or people. "
                "Do NOT infer, speculate, or add any information not present in the CONTEXT. "
                "Do NOT say 'Based on the context' — state the fact directly. "
                "If the context has no relevant information, say exactly: "
                "'The indexed documents do not contain information about this.'"
            )
            user_prompt = (
                f"Document: {source}\nContext:\n{_relevant_excerpt(context, query)}\n\n"
                f"Question: {query}\n\nAnswer:"
            )
            max_tokens = 500

        print("[FILE] Prompt created successfully")
        response = ollama.chat(
            model=model_name,
            options={"temperature": 0.0, "num_predict": max_tokens},
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user",   "content": user_prompt},
            ],
        )
        answer = response.get("message", {}).get("content", "").strip()
        if answer:
            return answer
    except Exception:
        pass
    return None


def handle_retrieval(
    query: str,
    vector_db,
    threshold: float,
    model_name: str,
    extra_dbs: list | None = None,
    last_file: Optional[str] = None,
    folder_path: Optional[str] = None,
) -> Tuple[Optional[str], Optional[str]]:
    """
    Answer a document query.

    Rule 1: If the user mentions a specific filename (any extension) →
            load ONLY that file and answer from it exclusively.

    Rule 2: Otherwise → use vector search across ALL indexed sources.
            When *extra_dbs* are provided (e.g. the Windows Documents store)
            results from all stores are merged before ranking so that every
            source contributes its best-matching chunk.

    Parameters
    ----------
    folder_path:
        When set by the access-control layer (ALLOW_FOLDER), all retrieval
        results are restricted to chunks whose ``metadata["source"]`` starts
        with this path.  Supersedes ``last_file`` restriction.
    """
    # Strip any [scope: ...] annotation that may have been injected by an older
    # orchestrator path — the folder is now passed explicitly via folder_path.
    query = re.sub(r'\s*\[scope:[^\]]*\]', '', query).strip()

    # ── Rule 0: query references an unauthorized path ──────────────────────────
    # If the user mentions any Windows path that is NOT under the authorized
    # documents root, return a clear "no access" message immediately.
    if _query_references_unauthorized_path(query):
        authorized = _get_authorized_docs_root()
        return (
            f"I cannot access information from that location. "
            f"I can only access documents inside {authorized}.",
            None,
        )

    # ── Rule 1: specific file mentioned ──────────────────────────────────────
    target_file = _detect_target_file(query)
    # File-folder security pre-check: when a folder scope is active the file
    # must live under that folder.  _detect_target_file only finds files in
    # DOCS_PATH (the project's local data/documents dir).  If DOCS_PATH is
    # outside folder_path we CANNOT conclude the file is absent — it may be
    # indexed in the vector store from folder_path.  In that case skip Rule 1
    # and let Rule 1b check the vector-store metadata instead.
    if target_file and folder_path:
        folder_norm = os.path.normcase(os.path.normpath(folder_path))
        docs_norm = os.path.normcase(os.path.normpath(DOCS_PATH))
        if not (docs_norm == folder_norm or docs_norm.startswith(folder_norm + os.sep)):
            log.debug(
                "Rule 1: DOCS_PATH %r is outside folder_path %r — "
                "deferring to Rule 1b vector-store check",
                DOCS_PATH, folder_path,
            )
            target_file = None  # fall through to Rule 1b below
    if target_file:
        fpath = os.path.join(DOCS_PATH, target_file)
        from services.academic_pdf_parser import AcademicPDFParser, handle_structured_academic_qa
        if AcademicPDFParser.is_academic_pdf(fpath):
            log.info("[STRUCTURED_ACADEMIC] Routing to specialized Academic Parser")
            ans = handle_structured_academic_qa(fpath, query)
            return ans, target_file

        ext = os.path.splitext(target_file)[1].lower()
        content, source = _load_file_content(target_file)

        # Image files: pass OCR text through LLM for explain/extract/summary queries;
        # fall back to returning the raw OCR text for bare "show me" requests.
        if ext in (".png", ".jpg", ".jpeg", ".webp"):
            if not content:
                return "I could not extract readable text from the image.", target_file
            if _is_summary_intent(query) or re.search(
                r'\b(explain|describe|extract|analyse|analyze|what|tell)\b',
                query, re.IGNORECASE,
            ):
                answer = _ask_llm(model_name, content, query, target_file,
                                  is_summary=_is_summary_intent(query))
                if answer:
                    return answer, target_file
            return f"OCR-extracted text from '{target_file}':\n{content}", target_file

        if content is None:
            return f"Could not read '{target_file}'.", target_file
        answer = _ask_llm(model_name, content, query, source,
                          is_summary=_is_summary_intent(query))
        if answer:
            return answer, source
        return content[:2000], source

    # ── Rule 1b: filename in query but file lives in a vector store ──────────────
    # _detect_target_file only inspects data/documents/ (project docs on disk).
    # When a filename pattern is detected in the query we do a metadata-filtered
    # search across ALL stores.  This rule is TERMINAL — if a filename is found
    # we never fall through to the broad vector search (Rule 2) so that unrelated
    # documents cannot mix into the context.
    filename_hint = _extract_filename_from_query(query)
    if filename_hint and not target_file:
        all_dbs_for_meta = [db for db in [vector_db] + list(extra_dbs or []) if db is not None]
        meta_docs = _search_by_filename_in_stores(filename_hint, all_dbs_for_meta)
        # Folder-scope security: when a folder is active keep only chunks from it.
        if meta_docs and folder_path:
            folder_fp = _normalize_path(folder_path)
            folder_basename = os.path.basename(os.path.normpath(folder_path)).lower()

            def _meta_in_folder(meta: dict) -> bool:
                src = _normalize_path(meta.get("source", ""))
                fname = _normalize_path(meta.get("file_name", ""))
                # file_path stores the full path; source/file_name are basenames only
                fpath_meta = _normalize_path(meta.get("file_path", ""))
                # Full normalized prefix match (file_path is the reliable key for new folders)
                if src.startswith(folder_fp) or fname.startswith(folder_fp) or fpath_meta.startswith(folder_fp):
                    return True
                # Folder basename contained in the stored path
                if folder_basename and (folder_basename in src or folder_basename in fname or folder_basename in fpath_meta):
                    return True
                return False

            meta_docs_filtered = [d for d in meta_docs if _meta_in_folder(d.metadata)]
            log.info(
                "handle_retrieval: Rule-1b folder filter: %d/%d chunk(s) matched for folder %r",
                len(meta_docs_filtered), len(meta_docs), folder_path,
            )
            if not meta_docs_filtered:
                # Strict: file found in stores but NOT in the active folder — do not leak it.
                log.info(
                    "handle_retrieval: Rule-1b folder filter found no match for %r under %r",
                    filename_hint, folder_path,
                )
                return (
                    f"The file '{filename_hint}' is not in '{folder_path}'.",
                    None,
                )
            meta_docs = meta_docs_filtered
        if meta_docs:
            resolved_source = _best_source_from_docs(meta_docs, filename_hint)
            log.info(
                "handle_retrieval: Rule-1b — %d chunk(s) from metadata search for %r (resolved source: %r)",
                len(meta_docs), filename_hint, resolved_source,
            )
            summary = _is_summary_intent(query)
            # Keyword relevance check: discard chunks with no query overlap unless
            # the user asked for a full summary/explanation (then we need all content)
            if not summary:
                relevant_chunks = _keyword_filter(meta_docs[:50], query)
                if not relevant_chunks:
                    log.info(
                        "handle_retrieval: Rule-1b keyword filter found no relevant chunks for %r",
                        query[:60],
                    )
                    return (
                        "I could not find this information in the document.",
                        resolved_source,
                    )
                chunks_for_context = relevant_chunks
            else:
                chunks_for_context = meta_docs[:50]
            # Use up to 50 chunks; deduplicate identical lines
            context = "\n\n".join(d.page_content for d in chunks_for_context)
            context = _deduplicate_lines(context)
            answer = _ask_llm(model_name, context, query, resolved_source, is_summary=summary)
            if answer:
                return answer, resolved_source
            return context[:4000], resolved_source
        else:
            # Image fallback: if the file exists in the authorized folder but has not
            # been indexed yet, try direct OCR so image explanation works immediately.
            if filename_hint.lower().endswith((".png", ".jpg", ".jpeg", ".webp")):
                authorized_root = _get_authorized_docs_root()
                img_path = _find_file_under_root(authorized_root, filename_hint)
                if img_path:
                    text = _ocr_image_file(img_path)
                    source = os.path.basename(img_path)
                    if text:
                        answer = _ask_llm(
                            model_name,
                            text,
                            query,
                            source,
                            is_summary=_is_summary_intent(query),
                        )
                        if answer:
                            return answer, source
                        return text[:4000], source
                    return (
                        "I found the image but could not extract readable text from it. "
                        "This assistant currently supports OCR-based image understanding. "
                        "If the image has little or no text, please upload a clearer text image."
                    ), source

            # Direct disk fallback: file exists in the granted/authorized folder but
            # has not been indexed yet.  Load and answer directly without indexing.
            _search_roots: list[str] = []
            if folder_path:
                _search_roots.append(folder_path)
            else:
                _search_roots.append(_get_authorized_docs_root())
            for _root in _search_roots:
                _disk_path = _find_file_under_root(_root, filename_hint)
                if _disk_path:
                    from services.academic_pdf_parser import AcademicPDFParser, handle_structured_academic_qa
                    if AcademicPDFParser.is_academic_pdf(_disk_path):
                        log.info("[STRUCTURED_ACADEMIC] Routing direct disk file to specialized Academic Parser")
                        ans = handle_structured_academic_qa(_disk_path, query)
                        return ans, os.path.basename(_disk_path)

                    log.info(
                        "handle_retrieval: Rule-1b direct disk load for %r from %r",
                        filename_hint, _root,
                    )
                    _content = _load_document_from_path(_disk_path)
                    if _content:
                        _fname_only = os.path.basename(_disk_path)
                        _answer = _ask_llm(
                            model_name, _content, query, _fname_only,
                            is_summary=_is_summary_intent(query),
                        )
                        if _answer:
                            return _answer, _fname_only
                        return _content[:4000], _fname_only
                    # File found on disk but couldn't be read
                    return (
                        f"Found '{filename_hint}' on disk but could not read its contents. "
                        "The file may be corrupted or use an unsupported format."
                    ), os.path.basename(_disk_path)

            # Filename detected but not found in any store or on disk — do NOT fall back to
            # vector search (which would return unrelated documents).  Return a
            # clear "not found" message so the user knows the file isn't indexed.
            log.warning(
                "handle_retrieval: Rule-1b — no indexed chunks for %r", filename_hint
            )
            if folder_path:
                # Check if the folder still has access permission
                try:
                    from core.permission_store import permission_store
                    if not permission_store.is_granted(folder_path):
                        return (
                            f"I no longer have access to {folder_path}.\n\n"
                            f"The folder permission may have been revoked. "
                            f"To grant access again, ask me to search or summarize "
                            f"files from that folder.",
                            None,
                        )
                except Exception:
                    pass
                return (
                    f"'{filename_hint}' was not found in '{folder_path}'.\n"
                    f"The folder is accessible — make sure the file exists there, "
                    f"then restart the assistant (or ask 'reindex {folder_path}') "
                    f"to trigger a fresh index scan.",
                    None,
                )
            authorized = _get_authorized_docs_root()
            return (
                f"'{filename_hint}' was not found in the indexed documents. "
                f"The only accessible folder is {authorized}. "
                "If the file exists there, try restarting the assistant "
                "to trigger a fresh index scan and then ask again.",
                None,
            )

    # ── Rule 2: vector search, multi-source aware ───────────────────────────────
    # Restrictions (mutually exclusive — folder_path takes precedence):
    #   _restrict_to_folder: normalised folder root; all results must live under it.
    #   _restrict_to_file:   basename of the active session document.
    _restrict_to_file: Optional[str] = None
    _restrict_to_folder: Optional[str] = None
    if folder_path:
        _restrict_to_file = None  # folder scope supersedes any file restriction
        _restrict_to_folder = _normalize_path(folder_path)
        log.info("handle_retrieval: Rule 2 restricted to folder=%r", folder_path)
        # Guard 1: folder must exist on disk
        if not os.path.isdir(folder_path):
            return (
                f"\u274c Invalid folder path\n\n"
                f"I could not find this folder:\n"
                f"\U0001f4c1 {folder_path}\n\n"
                "Please check the spelling, folder name, and path.",
                folder_path,
            )
        # Guard 2: folder must have at least one indexed file before running RAG
        # If no specific filename was mentioned and folder is not indexed,
        # try loading all files directly from disk instead of blocking.
        _filename_mentioned = _extract_filename_from_query(query)
        if not _filename_mentioned:
            try:
                from services.document_indexer_service import document_indexer_service as _dis_guard
                if not _dis_guard.is_folder_indexed(folder_path):
                    # Try to load all files from the folder directly
                    log.info("Folder not indexed, attempting direct folder load from %r", folder_path)
                    return _load_all_files_from_folder(folder_path, query, model_name)
            except Exception as _idx_guard_exc:
                log.debug("is_folder_indexed guard failed: %s — attempting direct load", _idx_guard_exc)
                # Fallback: try direct load even if guard fails
                return _load_all_files_from_folder(folder_path, query, model_name)
    elif last_file and _extract_filename_from_query(query) is None:
        _restrict_to_file = last_file
        log.info(
            "handle_retrieval: Rule 2 restricted to last_file=%r", last_file
        )

    all_dbs = [db for db in [vector_db] + list(extra_dbs or []) if db is not None]
    if not all_dbs:
        log.warning("[VECTOR_STORE] handle_retrieval: no vector stores available")
        # ── Fallback: load file directly from disk when last_file is known ──
        # This prevents silent failure when the orchestrator has no Chroma DB
        # but the active file path is resolvable on disk.
        _fallback_path: Optional[str] = None
        if last_file:
            # Try DOCS_PATH first, then current working directory
            _candidate = os.path.join(DOCS_PATH, last_file)
            if os.path.isfile(_candidate):
                _fallback_path = _candidate
            else:
                # last_file might be a full absolute path
                if os.path.isfile(last_file):
                    _fallback_path = last_file
        if _fallback_path:
            log.info(
                "[CHUNK_LOAD] No vector store — loading %r directly from disk",
                _fallback_path,
            )
            _fallback_content = _load_document_from_path(_fallback_path)
            if _fallback_content and not _fallback_content.startswith("Error"):
                _fname = os.path.basename(_fallback_path)
                log.info(
                    "[RETRIEVAL_READY] Direct fallback: %d chars from %r",
                    len(_fallback_content), _fname,
                )
                _is_sum = _is_summary_intent(query)
                _ans = _ask_llm(
                    model_name, _fallback_content, query, _fname, is_summary=_is_sum
                )
                if _ans:
                    return _ans, _fname
                return _fallback_content[:3000], _fname
        return None, None

    log.info(
        "handle_retrieval: querying %d vector store(s) for %r",
        len(all_dbs), query[:60],
    )
    for i, db in enumerate(all_dbs):
        label = getattr(db, "_persist_directory", None) or f"store[{i}]"
        log.debug("  store[%d]: %s", i, label)

    from engines.rag_engine import retrieve_top_k_multi
    print(f"[RETRIEVAL] Querying vector store for '{query}'")
    results = retrieve_top_k_multi(query, all_dbs, k=10)

    if not results:
        log.warning("handle_retrieval: no results from any vector store")
        return None, None

    # Keyword-boost re-sort: chunks containing query keywords bubble to the top
    # while vector similarity score acts as the tiebreaker within each tier.
    _kws = _extract_query_keywords(query)
    if _kws:
        results.sort(key=lambda item: (
            -sum(item[0].page_content.lower().count(kw) for kw in _kws),  # keyword hits desc
            item[1],  # then vector score asc (lower = more similar)
        ))
        log.debug(
            "handle_retrieval: keyword-boost re-sort applied (%d keywords: %s)",
            len(_kws), sorted(_kws),
        )

    # Folder filter — keep only results whose source path starts with the
    # authorised folder.  Supersedes the file-level restriction.
    if _restrict_to_folder:
        folder_basename_r2 = os.path.basename(os.path.normpath(folder_path)).lower()

        def _result_in_folder(doc) -> bool:
            src = _normalize_path(doc.metadata.get("source", ""))
            fname = _normalize_path(doc.metadata.get("file_name", ""))
            # file_path stores the full path; source/file_name are basenames only
            fpath_meta = _normalize_path(doc.metadata.get("file_path", ""))
            # Full normalized prefix match (file_path is the reliable key for new folders)
            if src.startswith(_restrict_to_folder) or fname.startswith(_restrict_to_folder) or fpath_meta.startswith(_restrict_to_folder):
                return True
            # Folder basename contained in the stored path
            if folder_basename_r2 and (folder_basename_r2 in src or folder_basename_r2 in fname or folder_basename_r2 in fpath_meta):
                return True
            return False

        strict_results = [(doc, score) for doc, score in results if _result_in_folder(doc)]
        log.info(
            "handle_retrieval: Rule 2 folder filter: %d/%d result(s) matched for folder %r",
            len(strict_results), len(results), folder_path,
        )
        if strict_results:
            log.debug(
                "handle_retrieval: Rule 2 folder filter sample — source=%r  file_path=%r",
                strict_results[0][0].metadata.get("source", ""),
                strict_results[0][0].metadata.get("file_path", ""),
            )
        if not strict_results:
            # Strict: no results in the requested folder — never fall back to other folders.
            log.info(
                "handle_retrieval: Rule 2 folder filter matched nothing for %r — no fallback",
                folder_path,
            )
            return (
                f"No relevant information found in:\n\U0001f4c1 {folder_path}\n\n"
                "The folder may not be indexed yet, or contains no searchable content.",
                folder_path,
            )
        results = strict_results

    if _restrict_to_file:
        fpath = os.path.join(DOCS_PATH, _restrict_to_file)
        if not os.path.exists(fpath) and folder_path:
            fpath = os.path.join(folder_path, _restrict_to_file)
            
        from services.academic_pdf_parser import AcademicPDFParser, handle_structured_academic_qa
        if os.path.exists(fpath) and AcademicPDFParser.is_academic_pdf(fpath):
            log.info("[STRUCTURED_ACADEMIC] Routing restricted file to specialized Academic Parser")
            ans = handle_structured_academic_qa(fpath, query)
            return ans, _restrict_to_file
            
        fname_lower = _restrict_to_file.lower()
        results = [
            (doc, score) for doc, score in results
            if doc.metadata.get("source", "").lower() == fname_lower
            or doc.metadata.get("file_name", "").lower() == fname_lower
        ]
        if not results:
            log.info(
                "handle_retrieval: no chunks matched last_file=%r after filter",
                _restrict_to_file,
            )
            return (
                f"No relevant information about this was found in '{_restrict_to_file}'.",
                _restrict_to_file,
            )
        print(f"[RETRIEVAL] Found {len(results)} initial chunks from {fname_lower}")

    # Include the best-scoring chunk from EVERY source document so that
    # smaller files (e.g. a 2-row CSV) are not crowded out by a large PDF
    # that dominates the top-k slots.
    best_per_source: dict = {}
    for doc, score in results:
        src = doc.metadata.get("source", "")
        if src not in best_per_source or score < best_per_source[src][1]:
            best_per_source[src] = (doc, score)

    # Sort sources by best score (lowest = most relevant first)
    sorted_sources = sorted(best_per_source.items(), key=lambda x: x[1][1])
    best_source = sorted_sources[0][0]

    # Only attribute sources whose best chunk is within relevance threshold
    relevant_sources = [(src, pair) for src, pair in sorted_sources if pair[1] <= threshold]
    if not relevant_sources:
        relevant_sources = sorted_sources[:1]  # always include at least the best match

    # One representative chunk per relevant source, then up to 4 extra from the
    # top source — this guarantees at least 5 chunks when a single source dominates.
    ordered_docs = [pair[0] for _, pair in relevant_sources]
    seen_ids = {id(d) for d in ordered_docs}
    extras = [
        doc for doc, _ in results
        if doc.metadata.get("source", "") == best_source and id(doc) not in seen_ids
    ][:4]
    all_docs = ordered_docs + extras
    print(f"[RERANK] Selected {len(all_docs)} top chunks")
    print(f"[CHUNK_COUNT] {len(all_docs)} chunks used for context")
    source = ", ".join(src for src, _ in relevant_sources)

    # Strict keyword relevance check — skip for summary/explain requests since
    # those intentionally need all document content, not query-specific passages.
    if not _is_summary_intent(query):
        kw_filtered = _keyword_filter(all_docs, query)
        if not kw_filtered:
            log.info(
                "handle_retrieval: Rule 2 keyword filter eliminated all chunks for %r",
                query[:60],
            )
            return "I could not find this information in the document.", source
        all_docs = kw_filtered

    context = "\n\n".join(doc.page_content for doc in all_docs)
    context = _deduplicate_lines(context)

    # For any CSV source in the results, replace its chunk with the full table so
    # the LLM sees ALL rows and cannot confuse different years as different companies.
    csv_supplements = []
    for src, _ in relevant_sources:
        if src.lower().endswith(".csv"):
            full_text, _ = _load_file_content(src)
            if full_text and not full_text.startswith("(Error"):
                csv_supplements.append(f"[Full table: {src}]\n{full_text}")
    if csv_supplements:
        context = "\n\n".join(csv_supplements) + "\n\n" + context

    answer = _ask_llm(model_name, context, query, source, is_summary=_is_summary_intent(query))
    if answer:
        return answer, source

    snippet = context if len(context) < 2000 else context[:2000] + "..."
    return snippet, source